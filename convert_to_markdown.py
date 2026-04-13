import os
import re
import csv
import sys
import json
import requests
import PIL.Image
from google import genai
from google.genai import types
import io
import time
from pathlib import Path
from datetime import datetime
from bs4 import BeautifulSoup
from markdownify import MarkdownConverter
from dotenv import load_dotenv
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
import openpyxl

# ---------------------------------------------------------------------------
# CONFIGURATION
# ---------------------------------------------------------------------------

load_dotenv()
GEMINI_API_KEY  = os.getenv("GEMINI_API_KEY")
VISION_MODEL    = "gemini-2.5-flash"
DESCRIBE_IMAGES  = True   # Set to False for a fast run without image descriptions
_gemini_client   = None   # initialised below once log file is open

RAW_DIR = Path("onenote_audit/01_Raw_Audit")
MD_DIR  = Path("onenote_audit/02_Markdown")

# ---------------------------------------------------------------------------
# LOGGING HELPER
# ---------------------------------------------------------------------------

def safe_log(msg):
    """Log to screen and file if the log file is open, otherwise just print."""
    print(msg)
    try:
        _log_file.write(msg + "\n")
    except Exception:
        pass  # log file may not be open yet during early startup


# ---------------------------------------------------------------------------
# MANIFEST HELPERS
# ---------------------------------------------------------------------------

_manifest_cache = {}

def _load_manifest(notebook_raw_dir):
    """Load and cache manifest.json for a notebook in 01_Raw_Audit."""
    key = str(notebook_raw_dir)
    if key not in _manifest_cache:
        path = Path(notebook_raw_dir) / "manifest.json"
        if path.exists():
            try:
                with open(path, "r", encoding="utf-8") as f:
                    _manifest_cache[key] = json.load(f)
            except Exception:
                _manifest_cache[key] = {"pages": {}}
        else:
            _manifest_cache[key] = {"pages": {}}
    return _manifest_cache[key]


def get_page_timestamp(html_file):
    """Return lastModifiedDateTime for a page from its notebook manifest."""
    rel   = html_file.parent.relative_to(RAW_DIR)
    parts = rel.parts
    if len(parts) < 2:
        return None
    manifest = _load_manifest(RAW_DIR / parts[0])
    page_key = "/".join(parts[1:])
    entry    = manifest.get("pages", {}).get(page_key)
    return entry.get("lastModifiedDateTime") if entry else None


def read_md_timestamp(md_file):
    """Read lastModifiedDateTime from an existing .md file's frontmatter."""
    try:
        text = md_file.read_text(encoding="utf-8")
        if text.startswith("---"):
            end = text.find("---", 3)
            if end != -1:
                for line in text[3:end].splitlines():
                    if line.strip().startswith("lastModifiedDateTime:"):
                        return line.split(":", 1)[1].strip().strip('"')
    except Exception:
        pass
    return None


# ---------------------------------------------------------------------------
# GEMINI VISION
# ---------------------------------------------------------------------------

def describe_image(image_path):
    """Send a local image to Gemini and return a plain-text description."""
    # Disable safety filters that can false-positive on technical diagrams,
    # network maps, security content, and screenshots of CLI/code output
    safety_settings = [
        types.SafetySetting(category='HARM_CATEGORY_HARASSMENT',        threshold='BLOCK_NONE'),
        types.SafetySetting(category='HARM_CATEGORY_HATE_SPEECH',        threshold='BLOCK_NONE'),
        types.SafetySetting(category='HARM_CATEGORY_SEXUALLY_EXPLICIT',  threshold='BLOCK_NONE'),
        types.SafetySetting(category='HARM_CATEGORY_DANGEROUS_CONTENT',  threshold='BLOCK_NONE'),
    ]
    try:
        img = PIL.Image.open(image_path)
        fmt = img.format or 'PNG'
        buf = io.BytesIO()

        # Gemini does not support TIFF — convert to PNG in memory
        if fmt == 'TIFF':
            img.convert('RGB').save(buf, format='PNG')
            mime_type = 'image/png'
        else:
            img.save(buf, format=fmt)
            mime_map  = {
                'JPEG': 'image/jpeg', 'PNG': 'image/png',
                'GIF':  'image/gif',  'WEBP': 'image/webp',
            }
            mime_type = mime_map.get(fmt, 'image/png')

        contents = [
            types.Part.from_bytes(data=buf.getvalue(), mime_type=mime_type),
            (
                "Describe this image in detail for use in a technical document. "
                "If it contains text, transcribe it exactly. "
                "If it shows a diagram, chart, network map, or technical drawing, "
                "explain precisely what it depicts. Be thorough."
            ),
        ]
        config = types.GenerateContentConfig(safety_settings=safety_settings)

        # Retry up to 3 times on transient errors (503 unavailable, 429 rate limit)
        max_retries = 3
        wait        = 10  # seconds — doubles each retry
        for attempt in range(1, max_retries + 1):
            try:
                response = _gemini_client.models.generate_content(
                    model=VISION_MODEL, contents=contents, config=config,
                )
                return response.text.strip()
            except Exception as e:
                err = str(e)
                transient = any(code in err for code in ("503", "429", "UNAVAILABLE"))
                if transient and attempt < max_retries:
                    safe_log(f"      [Retry {attempt}/{max_retries} in {wait}s]: {err}")
                    time.sleep(wait)
                    wait *= 2   # exponential backoff: 10s → 20s → 40s
                else:
                    safe_log(f"      [Image description error]: {e}")
                    return None
    except Exception as e:
        safe_log(f"      [Image description error]: {e}")
        return None


# ---------------------------------------------------------------------------
# ATTACHMENT TEXT EXTRACTION
# ---------------------------------------------------------------------------

def extract_attachment_text(file_path):
    """
    Extract readable text from a local attachment file.
    Supports: .docx, .pptx, .pdf, .xlsx, .csv, .txt
    Returns a markdown-formatted string of the content.
    """
    suffix = Path(file_path).suffix.lower()
    try:
        if suffix == ".docx":
            doc   = Document(file_path)
            lines = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(lines)

        elif suffix == ".pptx":
            prs   = Presentation(file_path)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"### Slide {i}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
            return "\n\n".join(lines)

        elif suffix == ".pdf":
            reader = PdfReader(file_path)
            lines  = []
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    lines.append(f"### Page {i}\n{text.strip()}")
            return "\n\n".join(lines)

        elif suffix == ".xlsx":
            wb    = openpyxl.load_workbook(file_path, data_only=True)
            lines = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                lines.append(f"### Sheet: {sheet_name}")
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(
                        str(cell) if cell is not None else "" for cell in row
                    )
                    if row_text.strip(" |"):
                        lines.append(row_text)
            return "\n\n".join(lines)

        elif suffix == ".csv":
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                rows = list(csv.reader(f))
            if not rows:
                return ""
            # Format as a markdown table
            header = "| " + " | ".join(rows[0]) + " |"
            sep    = "| " + " | ".join(["---"] * len(rows[0])) + " |"
            body   = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
            return "\n".join([header, sep, body])

        elif suffix == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read().strip()

        else:
            return None  # unsupported type — will fall back to reference only

    except Exception as e:
        safe_log(f"      [Extraction error for {Path(file_path).name}]: {e}")
        return None


# ---------------------------------------------------------------------------
# INDENTATION HELPER
# ---------------------------------------------------------------------------

def get_indent_level(style_str):
    """
    Parse a CSS style string and return an integer indent level.
    OneNote uses margin-left (in pt or px) to show nested indentation.
    Examples: margin-left:36pt → level 1,  margin-left:72pt → level 2
    """
    if not style_str:
        return 0
    match = re.search(r'margin-left\s*:\s*([\d.]+)(pt|px)', style_str)
    if not match:
        return 0
    value = float(match.group(1))
    unit  = match.group(2)
    base  = 36 if unit == "pt" else 40   # approx pixels per indent level
    return max(0, round(value / base))


# ---------------------------------------------------------------------------
# CUSTOM MARKDOWN CONVERTER
# ---------------------------------------------------------------------------

class OneNoteConverter(MarkdownConverter):
    """
    Extends markdownify to handle two OneNote-specific cases:
      1. <p style="margin-left:..."> — converts to indented markdown
      2. <img src="media/...">       — local image ref + optional Ollama description
    Everything else (headings, bold, italic, tables, lists, code) is handled
    by the base markdownify library.
    """

    def __init__(self, page_dir, **options):
        super().__init__(**options)
        self.page_dir = Path(page_dir)   # absolute path to the raw page folder

    # --- Indentation ---
    def convert_p(self, el, text, convert_as_inline=False, **kwargs):
        if not text.strip():
            return ""
        level  = get_indent_level(el.get("style", ""))
        indent = "    " * level           # 4 spaces per level
        # Apply indent to every line in case the paragraph wraps
        indented = "\n".join(
            indent + line if line.strip() else line
            for line in text.splitlines()
        )
        return f"{indented}\n\n"

    # --- Attachments (Word, PowerPoint, PDF, Excel, CSV, TXT) ---
    def convert_object(self, el, text, convert_as_inline=False, **kwargs):
        filename = el.get("data-attachment", "")
        data_src = el.get("data", "")

        if not filename:
            return ""

        # Resolve the local path — archive script saves to attachments/filename
        attach_path = (self.page_dir / data_src).resolve() if data_src else None

        header = f"\n---\n📎 **Attachment: {filename}**\n"

        if attach_path and attach_path.exists():
            content = extract_attachment_text(attach_path)
            if content:
                safe_log(f"      Extracting: {filename}")
                return f"{header}\n{content}\n\n---\n"
            else:
                # File type not supported for extraction — reference only
                return f"{header}*(Content extraction not supported for this file type)*\n\n---\n"
        else:
            # File wasn't downloaded or doesn't exist
            return f"{header}*(File not found in attachments folder)*\n\n---\n"

    # --- Images ---
    def convert_img(self, el, text, convert_as_inline=False, **kwargs):
        src = el.get("src", "")
        alt = el.get("alt", "Image")

        if not src or src.startswith("http"):
            return ""   # remote URLs were already handled in the archive step

        abs_path = (self.page_dir / src).resolve()
        img_ref  = f"\n![{alt}]({abs_path})\n"

        if DESCRIBE_IMAGES and abs_path.exists():
            safe_log(f"      Describing: {src} ...")
            description = describe_image(abs_path)
            if description:
                quoted = "\n".join(f"> {line}" for line in description.splitlines())
                return f"{img_ref}\n> **[Image Description]**\n{quoted}\n\n"

        return f"{img_ref}\n"


# ---------------------------------------------------------------------------
# PAGE CONVERTER
# ---------------------------------------------------------------------------

def convert_page(html_file, page_dir, md_file, metadata):
    """Read index.html, convert to Markdown with YAML frontmatter, write .md file."""

    with open(html_file, "r", encoding="utf-8") as f:
        html = f.read()

    soup = BeautifulSoup(html, "html.parser")

    # Pull creation date from OneNote's <meta name="created"> if present
    meta_tag = soup.find("meta", attrs={"name": "created"})
    if meta_tag:
        metadata["created"] = meta_tag.get("content", "")

    # --- YAML frontmatter ---
    fm_lines = ["---"]
    for key, val in metadata.items():
        if val:
            # Escape any double quotes inside values
            safe_val = str(val).replace('"', '\\"')
            fm_lines.append(f'{key}: "{safe_val}"')
    fm_lines.append("---\n")
    frontmatter = "\n".join(fm_lines)

    # --- Convert HTML body to Markdown ---
    body      = soup.find("body") or soup
    converter = OneNoteConverter(page_dir=page_dir, heading_style="ATX")
    md_body   = converter.convert(str(body))

    # Collapse 3+ consecutive blank lines down to 2
    md_body = re.sub(r"\n{3,}", "\n\n", md_body).strip()

    # --- Write output ---
    os.makedirs(md_file.parent, exist_ok=True)
    with open(md_file, "w", encoding="utf-8") as f:
        f.write(frontmatter + "\n" + md_body + "\n")


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------

print("=" * 60)
print("  OneNote HTML → Markdown Converter")
print("=" * 60)
print(f"  Source : {RAW_DIR.resolve()}")
print(f"  Output : {MD_DIR.resolve()}")
print(f"  Model  : {VISION_MODEL} (Gemini)")
print(f"  Images : {'Describe via Gemini' if DESCRIBE_IMAGES else 'Reference only (fast mode)'}")
print()

# --- Log file setup ---
os.makedirs("onenote_audit", exist_ok=True)
log_path  = os.path.join(
    "onenote_audit",
    f"02_Markdown-log_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}.txt"
)
_log_file = open(log_path, "w", encoding="utf-8", buffering=1)
print(f"Logging to: {log_path}\n")

def log(msg=""):
    """Print to screen and write to log file immediately."""
    print(msg)
    _log_file.write(msg + "\n")

# --- Initialise Gemini ---
if DESCRIBE_IMAGES:
    if not GEMINI_API_KEY:
        log("  WARNING: GEMINI_API_KEY not found in .env. Disabling image descriptions.\n")
        DESCRIBE_IMAGES = False
    else:
        _gemini_client = genai.Client(api_key=GEMINI_API_KEY)
        log(f"  Gemini ready. Model: {VISION_MODEL}\n")

# --- Walk and convert ---
total_converted = 0
total_updated   = 0
total_skipped   = 0
total_errors    = 0

all_pages = sorted(RAW_DIR.rglob("index.html"))

# --- Notebook selection ---
# Derive notebook names from top-level folders in 01_Raw_Audit
notebooks = sorted([d.name for d in RAW_DIR.iterdir() if d.is_dir()])
print("Available notebooks:")
for i, nb in enumerate(notebooks, 1):
    print(f"  {i:>2}. {nb}")

print()
print("  a           — convert all notebooks")
print("  s <numbers> — skip these notebooks       e.g. s 1 3 5")
print("  p <numbers> — pull only these notebooks  e.g. p 2 4")
print()

selection = input("Your choice: ").strip().lower()

if selection == "a" or selection == "":
    print(f"\nConverting all {len(notebooks)} notebooks.\n")
elif selection.startswith("s "):
    skip_nums = set(int(x) for x in selection[2:].split() if x.isdigit())
    skip_names = {notebooks[i-1] for i in skip_nums if 1 <= i <= len(notebooks)}
    all_pages = [p for p in all_pages if p.relative_to(RAW_DIR).parts[0] not in skip_names]
    print(f"\nSkipping {len(skip_nums)} notebook(s). Processing remaining pages.\n")
elif selection.startswith("p "):
    pull_nums = set(int(x) for x in selection[2:].split() if x.isdigit())
    pull_names = {notebooks[i-1] for i in pull_nums if 1 <= i <= len(notebooks)}
    all_pages = [p for p in all_pages if p.relative_to(RAW_DIR).parts[0] in pull_names]
    print(f"\nConverting {len(pull_nums)} notebook(s).\n")
else:
    print(f"\nUnrecognised input — converting all notebooks.\n")

# Log the selection
log(f"Notebooks selected for this run:")
selected_notebooks = sorted(set(p.relative_to(RAW_DIR).parts[0] for p in all_pages))
for nb in selected_notebooks:
    log(f"  - {nb}")
log()

log(f"Found {len(all_pages)} pages to process.\n")

for html_file in all_pages:
    page_dir = html_file.parent
    rel_path = page_dir.relative_to(RAW_DIR)
    parts    = rel_path.parts

    # Extract metadata from folder path
    # Structure: Notebook / [SectionGroup /] Section / PageTitle
    notebook   = parts[0] if len(parts) >= 1 else "Unknown"
    page_title = parts[-1] if len(parts) >= 1 else "Unknown"
    section    = " / ".join(parts[1:-1]) if len(parts) >= 3 else parts[1] if len(parts) == 2 else ""

    manifest_ts = get_page_timestamp(html_file)

    metadata = {
        "notebook"             : notebook,
        "section"              : section,
        "title"                : page_title,
        "created"              : "",   # filled in by convert_page if available
        "lastModifiedDateTime" : manifest_ts or "",
        "source"               : str(html_file.resolve()),
    }

    # Output: 02_Markdown/Notebook/Section/PageTitle.md  (flat .md, no subfolder)
    md_file = MD_DIR / rel_path.parent / f"{page_title}.md"

    # --- Resume: skip unless page has been updated since last conversion ---
    is_update = False
    if md_file.exists():
        if manifest_ts and read_md_timestamp(md_file) == manifest_ts:
            log(f"  [Skip] {notebook} > {section} > {page_title}")
            total_skipped += 1
            continue
        elif manifest_ts:
            # Timestamp changed — re-convert
            is_update = True
            log(f"  [Update] {notebook} > {section} > {page_title}")
        else:
            # No manifest — fall back to old behaviour (skip if file exists)
            log(f"  [Skip] {notebook} > {section} > {page_title}")
            total_skipped += 1
            continue
    else:
        log(f"  Converting: {notebook} > {section} > {page_title}")

    try:
        convert_page(html_file, page_dir, md_file, metadata)
        if is_update:
            total_updated += 1
        else:
            total_converted += 1
    except Exception as e:
        log(f"    [ERROR]: {e}")
        total_errors += 1

log(f"\n{'=' * 60}")
log(f"  CONVERSION COMPLETE")
log(f"{'=' * 60}")
log(f"  Converted : {total_converted}")
log(f"  Updated   : {total_updated}")
log(f"  Skipped   : {total_skipped}  (unchanged)")
log(f"  Errors    : {total_errors}")
log(f"  Output    : {MD_DIR.resolve()}")
_log_file.close()
